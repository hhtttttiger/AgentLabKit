import io
from pptx import Presentation
from retrieval.engines.local_engine.document_loaders.base_loader import ServiceLoader
from retrieval.engines.local_engine.document_loaders.base import Document, Page

class PPTLoader(ServiceLoader):
    def load(self, file_bytes: bytes, file_info, **kwargs) -> Document:
        stream = io.BytesIO(file_bytes)
        prs = Presentation(stream)
        
        content = ""
        page_list = []
        
        for i, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
            
            text = " ".join(text_runs)
            if text:
                content += text + "\n"
                page_list.append(Page(content=text, page_number=i+1))
                
        return Document(content=content, page_list=page_list, is_md=False)
